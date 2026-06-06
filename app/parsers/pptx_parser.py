from pptx import Presentation
import io

def parse_pptx(file_bytes):
    """
    Parses a PowerPoint file stream from memory and extracts text into Markdown.
    """
    markdown_lines = []
    
    # Open the PowerPoint from the byte stream
    prs = Presentation(io.BytesIO(file_bytes))
    
    for i, slide in enumerate(prs.slides):
        markdown_lines.append(f"## Slide {i + 1}\n")
        
        # Look at every shape/element on the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Extract and clean the text
                cleaned_text = shape.text.strip().replace('\v', '\n') 
                markdown_lines.append(cleaned_text)
                markdown_lines.append("\n")
                
        # Add a visual separator between slides
        markdown_lines.append("---\n")
        
    return "\n".join(markdown_lines)